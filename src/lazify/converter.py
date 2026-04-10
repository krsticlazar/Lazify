from __future__ import annotations

import subprocess
import sys
import tempfile
from pathlib import Path

import mammoth
from pdfminer.high_level import extract_text
from pptx import Presentation


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".ppt",
    ".pptx",
}


class ConverterError(Exception):
    pass


class UnsupportedFileTypeError(ConverterError):
    pass


class FileConverter:
    def is_supported(self, source_path: str | Path) -> bool:
        return Path(source_path).suffix.lower() in SUPPORTED_EXTENSIONS

    def default_output_path(
        self,
        source_path: str | Path,
        output_directory: str | Path | None = None,
    ) -> Path:
        source = Path(source_path).expanduser().resolve()
        if output_directory:
            target_directory = Path(output_directory).expanduser().resolve()
            return target_directory / f"{source.stem}.md"
        return source.with_suffix(".md")

    def suggest_output_path(
        self,
        source_path: str | Path,
        reserved_paths: set[Path] | None = None,
        output_directory: str | Path | None = None,
    ) -> Path:
        candidate = self.default_output_path(source_path, output_directory=output_directory)
        reserved = {Path(path).resolve() for path in (reserved_paths or set())}

        if candidate.resolve() not in reserved and not candidate.exists():
            return candidate

        counter = 1
        while True:
            fallback = candidate.with_name(f"{candidate.stem} ({counter}){candidate.suffix}")
            resolved_fallback = fallback.resolve()
            if resolved_fallback not in reserved and not fallback.exists():
                return fallback
            counter += 1

    def convert_file(self, source_path: str | Path) -> str:
        path = Path(source_path).expanduser().resolve()
        self._validate_source(path)

        try:
            if path.suffix.lower() == ".pdf":
                markdown_text = self._convert_pdf(path)
            elif path.suffix.lower() == ".docx":
                markdown_text = self._convert_docx(path)
            elif path.suffix.lower() == ".ppt":
                markdown_text = self._convert_ppt(path)
            elif path.suffix.lower() == ".pptx":
                markdown_text = self._convert_pptx(path)
            else:
                raise UnsupportedFileTypeError(f"Unsupported file type: {path.suffix or 'unknown'}")
        except ConverterError:
            raise
        except Exception as exc:
            raise ConverterError(self._format_exception(path, exc)) from exc

        if not markdown_text.strip():
            raise ConverterError(f"{path.name}: No readable content found.")

        return markdown_text.strip()

    def save_markdown(
        self,
        source_path: str | Path,
        markdown_text: str,
        target_path: str | Path | None = None,
    ) -> Path:
        source = Path(source_path).expanduser().resolve()
        destination = Path(target_path).expanduser().resolve() if target_path else self.default_output_path(source)

        try:
            destination.parent.mkdir(parents=True, exist_ok=True)
            destination.write_text(markdown_text, encoding="utf-8")
        except OSError as exc:
            raise ConverterError(self._format_exception(destination, exc)) from exc

        return destination

    def _validate_source(self, path: Path) -> None:
        if not path.exists():
            raise ConverterError(f"File not found: {path.name}")

        if not path.is_file():
            raise ConverterError(f"Not a file: {path.name}")

        if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            raise UnsupportedFileTypeError(f"Unsupported file type: {path.suffix or 'unknown'}")

    def _convert_pdf(self, path: Path) -> str:
        text = extract_text(str(path))
        return self._normalize_text(text)

    def _convert_docx(self, path: Path) -> str:
        with path.open("rb") as handle:
            result = mammoth.convert_to_markdown(handle)
        return self._normalize_text(result.value)

    def _convert_ppt(self, path: Path) -> str:
        if sys.platform != "win32":
            raise ConverterError(f"{path.name}: Legacy .ppt conversion is only available on Windows.")

        with tempfile.TemporaryDirectory(prefix="lazify-ppt-") as temp_dir:
            converted_path = Path(temp_dir) / f"{path.stem}.pptx"
            self._export_ppt_to_pptx(path, converted_path)
            return self._convert_pptx(converted_path)

    def _convert_pptx(self, path: Path) -> str:
        presentation = Presentation(str(path))
        sections: list[str] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_blocks: list[str] = [f"# Slide {slide_index}"]

            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                text = self._normalize_text(shape.text)
                if not text:
                    continue

                if getattr(shape, "has_text_frame", False) and shape == getattr(slide.shapes, "title", None):
                    slide_blocks.append(f"## {text}")
                    continue

                paragraphs = [line.strip() for line in text.splitlines() if line.strip()]
                for paragraph in paragraphs:
                    slide_blocks.append(f"- {paragraph}")

            sections.append("\n\n".join(slide_blocks))

        return "\n\n---\n\n".join(section for section in sections if section.strip())

    def _export_ppt_to_pptx(self, source_path: Path, target_path: Path) -> None:
        script = f"""
$ErrorActionPreference = 'Stop'
$source = {self._powershell_literal(source_path)}
$target = {self._powershell_literal(target_path)}
$powerpoint = $null
$presentation = $null

try {{
    $powerpoint = New-Object -ComObject PowerPoint.Application
    $presentation = $powerpoint.Presentations.Open($source, $true, $false, $false)
    $presentation.SaveAs($target, 24)
}}
finally {{
    if ($presentation -ne $null) {{
        $presentation.Close()
        [void][System.Runtime.InteropServices.Marshal]::ReleaseComObject($presentation)
    }}
    if ($powerpoint -ne $null) {{
        $powerpoint.Quit()
        [void][System.Runtime.InteropServices.Marshal]::ReleaseComObject($powerpoint)
    }}
    [System.GC]::Collect()
    [System.GC]::WaitForPendingFinalizers()
}}
"""

        try:
            result = subprocess.run(
                [
                    "powershell.exe",
                    "-NoProfile",
                    "-NonInteractive",
                    "-ExecutionPolicy",
                    "Bypass",
                    "-Command",
                    script,
                ],
                capture_output=True,
                text=True,
                check=False,
                timeout=180,
            )
        except OSError as exc:
            raise ConverterError(f"{source_path.name}: Unable to launch PowerShell for legacy .ppt conversion.") from exc
        except subprocess.TimeoutExpired as exc:
            raise ConverterError(f"{source_path.name}: Timed out while converting the legacy .ppt file.") from exc

        if result.returncode != 0:
            details = "\n".join(part.strip() for part in (result.stderr, result.stdout) if part and part.strip())
            lowered = details.lower()
            if "powerpoint.application" in lowered or "class not registered" in lowered:
                raise ConverterError(
                    f"{source_path.name}: Legacy .ppt conversion requires Microsoft PowerPoint to be installed."
                )

            if details:
                raise ConverterError(f"{source_path.name}: {details.splitlines()[0]}")
            raise ConverterError(f"{source_path.name}: Legacy .ppt conversion failed.")

        if not target_path.exists():
            raise ConverterError(f"{source_path.name}: PowerPoint did not produce a .pptx export.")

    @staticmethod
    def _normalize_text(text: str) -> str:
        lines = [line.rstrip() for line in text.replace("\r\n", "\n").split("\n")]
        compact: list[str] = []
        previous_blank = False

        for line in lines:
            stripped = line.strip()
            if not stripped:
                if not previous_blank:
                    compact.append("")
                previous_blank = True
                continue

            compact.append(stripped)
            previous_blank = False

        return "\n".join(compact).strip()

    @staticmethod
    def _format_exception(path: Path, exc: Exception) -> str:
        message = str(exc).strip()
        if not message:
            message = exc.__class__.__name__
        return f"{path.name}: {message}"

    @staticmethod
    def _powershell_literal(path: Path) -> str:
        return "'" + str(path).replace("'", "''") + "'"
